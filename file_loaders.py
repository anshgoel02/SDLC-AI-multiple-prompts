from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from pypdf import PdfReader

from .config import SUPPORTED_EXTENSIONS

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from docx import Document
except Exception:
    Document = None


def read_pdf_text(path: Path) -> str:
    reader = PdfReader(str(path))
    parts = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    return "\n".join(parts)


def read_pptx_text(path: Path) -> str:
    if Presentation is None:
        return ""
    pres = Presentation(str(path))
    parts = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return "\n".join(parts)


def read_docx_text(path: Path) -> str:
    if Document is None:
        return ""
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def load_source_text(path: Path) -> str:
    if path.suffix.lower() == ".pdf":
        return read_pdf_text(path)
    if path.suffix.lower() == ".pptx":
        return read_pptx_text(path)
    if path.suffix.lower() == ".docx":
        return read_docx_text(path)
    return path.read_text(encoding="utf-8", errors="ignore")


def iter_input_files(paths: Sequence[str]) -> Iterable[Path]:
    for raw in paths:
        path = Path(raw)
        if not path.exists():
            continue
        if path.is_file():
            if path.suffix.lower() in SUPPORTED_EXTENSIONS:
                yield path
        else:
            for file_path in path.rglob("*"):
                if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
                    yield file_path
